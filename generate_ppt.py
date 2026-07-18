from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_slide_background(slide, prs):
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    bg_shape = slide.shapes.add_shape(
        1, left, top, width, height # 1 is MSO_SHAPE.RECTANGLE
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(15, 23, 42) # Tailwind slate-900
    bg_shape.line.fill.background() # No line
    
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)

def style_text_frame(tf, is_title=False):
    for p in tf.paragraphs:
        p.font.name = 'Segoe UI'
        if is_title:
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255) # White
        else:
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(203, 213, 225) # slate-300

def add_logo(slide, logo_path, small=False):
    if not os.path.exists(logo_path): return
    if small:
        slide.shapes.add_picture(logo_path, Inches(8.5), Inches(0.2), width=Inches(1.0))
    else:
        slide.shapes.add_picture(logo_path, Inches(4), Inches(0.5), width=Inches(2.0))


def add_title_slide(prs, title, subtitle, logo_path):
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    style_text_frame(title_shape.text_frame, is_title=True)
    style_text_frame(subtitle_shape.text_frame, is_title=False)
    
    add_logo(slide, logo_path, small=False)

def add_content_slide(prs, title, bullet_points, logo_path):
    slide_layout = prs.slide_layouts[1] # Bullet slide
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, prs)
    
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    style_text_frame(title_shape.text_frame, is_title=True)
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        
    style_text_frame(tf, is_title=False)
    add_logo(slide, logo_path, small=True)

def main():
    prs = Presentation()
    logo_path = r"c:\Users\LC\Desktop\indisearch\frontend\public\logo.png"
    
    add_title_slide(prs, "INDISearch", "A Seamless, Intelligent, and Multi-Lingual Search Experience", logo_path)
    
    add_content_slide(prs, "INTRODUCTION", [
        "INDISearch is a modern, intelligent, and multi-lingual search engine.",
        "Designed to deliver fast, accurate, and visually appealing search results.",
        "Focuses on a clean, intuitive, and professional user experience."
    ], logo_path)

    add_content_slide(prs, "NEED FOR PROJECT", [
        "To overcome language barriers in accessing global information.",
        "Existing search engines may lack quick, localized AI-driven summaries.",
        "A requirement for a visually polished and responsive search platform."
    ], logo_path)

    add_content_slide(prs, "PROBLEM STATEMENT", [
        "Users often struggle to find accurate, summarized information quickly in their native language.",
        "Information overload makes it difficult to extract key insights effortlessly.",
        "Traditional search interfaces can feel cluttered and unintuitive."
    ], logo_path)

    add_content_slide(prs, "OBJECTIVES", [
        "To build a real-time search engine that dynamically displays user-queried content.",
        "Implement AI Overviews for intelligent query summaries.",
        "Integrate a robust Multi-Language Translation System.",
        "Provide Knowledge Panels for quick access to entity information."
    ], logo_path)

    add_content_slide(prs, "LANGUAGE USED", [
        "Frontend: React.js with Vite for fast rendering and development.",
        "Backend: Node.js for scalable API and data handling.",
        "Styling: Modern CSS frameworks for a responsive UI.",
        "Deployment: Vercel for fast, global edge delivery."
    ], logo_path)

    add_content_slide(prs, "RESULT", [
        "A fully functional, responsive search engine with a familiar layout.",
        "Successful integration of AI-driven quick summaries and entity panels.",
        "Seamless translation features enabling global accessibility."
    ], logo_path)

    add_content_slide(prs, "CHALLENGES", [
        "Managing real-time data fetching and ensuring low latency.",
        "Designing a responsive layout that works flawlessly across all devices.",
        "Integrating and optimizing AI models for accurate summaries."
    ], logo_path)

    add_content_slide(prs, "FUTURE WORK", [
        "Implementing personalized search results based on user preferences.",
        "Expanding language support and refining translation accuracy.",
        "Adding voice search and image-based search capabilities."
    ], logo_path)
    
    add_content_slide(prs, "CONCLUSION", [
        "INDISearch redefines how users interact with search engines by combining AI and intuitive design.",
        "It provides a powerful foundation for the future of multi-lingual search technology.",
        "A step forward in making global knowledge accessible to everyone."
    ], logo_path)
    
    output_path = r"c:\Users\LC\Desktop\indisearch\INDISearch_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    main()
